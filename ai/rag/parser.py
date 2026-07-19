"""文档解析器：支持多种常见文件格式"""
import io
import re
from pathlib import Path
from typing import List

from pypdf import PdfReader
from docx import Document


class DocumentParser:
    """统一的文档解析器
    
    支持的文件格式：
    - PDF (.pdf)
    - Word (.docx, .doc)
    - PowerPoint (.pptx)
    - Excel (.xlsx, .xls)
    - 文本文件 (.txt, .md, .csv)
    - 网页文件 (.html, .htm)
    - RTF (.rtf)
    """
    
    @staticmethod
    def parse_pdf(file_bytes: bytes) -> str:
        """解析 PDF 文件，返回纯文本"""
        text_parts = []
        try:
            reader = PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"PDF 解析失败: {str(e)}")
    
    @staticmethod
    def parse_word(file_bytes: bytes) -> str:
        """解析 Word 文件（.docx），返回纯文本"""
        try:
            doc = Document(io.BytesIO(file_bytes))
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            raise ValueError(f"Word 解析失败: {str(e)}")
    
    @staticmethod
    def parse_text(file_bytes: bytes, encoding: str = 'utf-8') -> str:
        """解析纯文本文件（.txt, .md, .csv等）"""
        # 尝试多种编码
        encodings = [encoding, 'utf-8', 'gbk', 'gb2312', 'gb18030', 'latin-1']
        
        for enc in encodings:
            try:
                return file_bytes.decode(enc)
            except UnicodeDecodeError:
                continue
        
        raise ValueError("无法解码文本文件，请确保文件编码为 UTF-8 或 GBK")
    
    @staticmethod
    def parse_html(file_bytes: bytes) -> str:
        """解析 HTML 文件，去除标签返回纯文本"""
        try:
            text = file_bytes.decode('utf-8')
            # 去除 HTML 标签
            text = re.sub(r'<[^>]+>', ' ', text)
            # 去除多余空白
            text = re.sub(r'\s+', ' ', text)
            # 还原段落
            text = re.sub(r'\s*([。！？.!?])\s*', r'\1\n', text)
            return text.strip()
        except Exception as e:
            raise ValueError(f"HTML 解析失败: {str(e)}")
    
    @staticmethod
    def parse_pptx(file_bytes: bytes) -> str:
        """解析 PowerPoint 文件（.pptx）"""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            text_parts = []
            
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    text_parts.append("\n".join(slide_text))
            
            return "\n\n=== 新幻灯片 ===\n\n".join(text_parts)
        except ImportError:
            raise ValueError("解析 PPT 需要安装 python-pptx: pip install python-pptx")
        except Exception as e:
            raise ValueError(f"PPT 解析失败: {str(e)}")

    @staticmethod
    def _pptx_picture_blobs_from_shape(shape, acc: List[bytes], max_n: int, max_bytes: int) -> None:
        """递归收集幻灯片中的嵌入图片（含组合图形内）。"""
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            return
        try:
            st = shape.shape_type
        except Exception:
            return
        if st == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                if blob and len(blob) <= max_bytes and len(acc) < max_n:
                    acc.append(blob)
            except Exception:
                pass
        elif st == MSO_SHAPE_TYPE.GROUP:
            for sub in getattr(shape, "shapes", []):
                DocumentParser._pptx_picture_blobs_from_shape(sub, acc, max_n, max_bytes)

    @staticmethod
    def parse_pptx_slides(file_bytes: bytes) -> List[dict]:
        """按页解析 pptx：文本 + 每页嵌入图片二进制（供 VLM 使用）。

        返回元素字段：
        - slide_no, title, raw_text, image_blobs（内部使用，不出现在最终 API 时可剥离）
        - visual_elements 占位为空列表，由 PreclassPptLLM+VLM 填充
        """
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            slides: List[dict] = []
            max_images = 3
            max_bytes = 1_500_000  # 单张约 1.5MB，避免超出视觉 API 限制
            for idx, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                raw_text = "\n".join(texts)
                title = texts[0][:120] if texts else ""
                blobs: List[bytes] = []
                for shape in slide.shapes:
                    DocumentParser._pptx_picture_blobs_from_shape(shape, blobs, max_images, max_bytes)
                slides.append(
                    {
                        "slide_no": idx,
                        "title": title,
                        "raw_text": raw_text,
                        "image_blobs": blobs,
                        "visual_elements": [],
                    }
                )
            return slides
        except ImportError:
            raise ValueError("解析 PPT 需要安装 python-pptx: pip install python-pptx")
        except Exception as e:
            raise ValueError(f"PPT 解析失败: {str(e)}")
    
    @staticmethod
    def parse_excel(file_bytes: bytes) -> str:
        """解析 Excel 文件（.xlsx, .xls）"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"【工作表: {sheet_name}】"]
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) for cell in row if cell is not None]
                    if row_text:
                        sheet_text.append(" | ".join(row_text))
                
                if len(sheet_text) > 1:
                    text_parts.append("\n".join(sheet_text))
            
            return "\n\n".join(text_parts)
        except ImportError:
            raise ValueError("解析 Excel 需要安装 openpyxl: pip install openpyxl")
        except Exception as e:
            raise ValueError(f"Excel 解析失败: {str(e)}")
    
    @staticmethod
    def parse_rtf(file_bytes: bytes) -> str:
        """解析 RTF 文件，提取纯文本"""
        try:
            text = file_bytes.decode('utf-8', errors='ignore')
            # 简单的 RTF 标签去除
            # 去除 RTF 控制字
            text = re.sub(r'\\[a-z]+\d*\s?', ' ', text)
            # 去除 RTF 花括号
            text = re.sub(r'[{}]', '', text)
            # 去除 RTF 头部
            text = re.sub(r'\\rtf1\s*', '', text)
            # 处理 Unicode 转义
            text = re.sub(r"\\'([0-9a-fA-F]{2})", lambda m: chr(int(m.group(1), 16)), text)
            # 去除多余空白
            text = re.sub(r'\s+', ' ', text)
            return text.strip()
        except Exception as e:
            raise ValueError(f"RTF 解析失败: {str(e)}")
    
    @classmethod
    def parse(cls, file_bytes: bytes, filename: str) -> str:
        """
        根据文件名自动选择解析器
        
        Args:
            file_bytes: 文件二进制内容
            filename: 文件名（用于判断类型）
        
        Returns:
            提取的纯文本
        
        Raises:
            ValueError: 不支持的文件类型或解析失败
        """
        ext = Path(filename).suffix.lower()
        
        # PDF
        if ext == '.pdf':
            return cls.parse_pdf(file_bytes)
        
        # Word
        elif ext in ['.docx']:
            return cls.parse_word(file_bytes)
        elif ext == '.doc':
            # .doc 是旧格式，尝试解析或提示转换
            raise ValueError(".doc 格式较旧，请转换为 .docx 格式后上传")
        
        # PowerPoint
        elif ext == '.pptx':
            return cls.parse_pptx(file_bytes)
        elif ext in ['.ppt']:
            raise ValueError(".ppt 格式较旧，请转换为 .pptx 格式后上传")
        
        # Excel
        elif ext in ['.xlsx', '.xls']:
            return cls.parse_excel(file_bytes)
        
        # 文本文件
        elif ext in ['.txt', '.md', '.markdown', '.csv', '.json', '.xml', '.py', '.js', '.java', '.c', '.cpp', '.h']:
            return cls.parse_text(file_bytes)
        
        # HTML
        elif ext in ['.html', '.htm']:
            return cls.parse_html(file_bytes)
        
        # RTF
        elif ext == '.rtf':
            return cls.parse_rtf(file_bytes)
        
        # 其他：尝试作为文本读取
        else:
            try:
                return cls.parse_text(file_bytes)
            except UnicodeDecodeError:
                raise ValueError(
                    f"不支持的文件类型: {ext}\n"
                    f"支持的格式：.pdf, .docx, .pptx, .xlsx, .xls, .txt, .md, .csv, .html, .rtf 等\n"
                    f"旧版 .doc 和 .ppt 请转换为新版格式后上传"
                )
    
    @classmethod
    def get_supported_extensions(cls) -> list:
        """获取支持的文件扩展名列表"""
        return [
            '.pdf', '.docx', '.pptx', '.xlsx', '.xls',
            '.txt', '.md', '.markdown', '.csv',
            '.html', '.htm', '.rtf',
            '.json', '.xml',
            '.py', '.js', '.java', '.c', '.cpp', '.h'
        ]
